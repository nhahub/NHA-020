from deep_translator import GoogleTranslator
from pypdf import PdfReader
from pptx import Presentation
import json
from langchain.text_splitter import RecursiveCharacterTextSplitter
from transformers import pipeline
from langchain.llms import HuggingFacePipeline
from langchain import PromptTemplate
import os
from langchain.vectorstores import FAISS
from langchain.chains import RetrievalQA
from langchain_openai import ChatOpenAI
from langchain_community.embeddings import HuggingFaceEmbeddings

os.environ["OPENAI_API_KEY"] =#this is Api key 
def read_pdf(file_path):
    reader = PdfReader(file_path)
    text = ''
    for page_num in range(len(reader.pages)):
        text += reader.pages[page_num].extract_text()
    return text

def read_powerpoint(file_path):
    presentation = Presentation(file_path)
    slides_text = []
    for slide in presentation.slides:
        slide_text = []
        for shape in slide.shapes:
            try:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            slide_text.append(run.text)
            except AttributeError:
                pass
        slides_text.append(' '.join(slide_text))
    return str(slides_text)

def read_notebook(notebook_path):
    with open(notebook_path, 'r', encoding='utf-8') as f:
        notebook = json.load(f)
    code_cells = []
    for cell in notebook['cells']:
        if cell['cell_type'] == 'code':
            code = ''.join(cell['source'])
            code_cells.append(code)
    all_code = '\n'.join(code_cells)
    return str(all_code)

def read_py(py_file_path):
    with open(py_file_path, 'r', encoding='utf-8') as py_file:
        content = py_file.read()
    return str(content)

def return_content(file_path):
    file_name, file_extension = os.path.splitext(file_path)
    file_extension = file_extension.lower()
    text = ''
    if file_extension == '.pdf':
        text = read_pdf(file_path)
    elif file_extension == '.pptx':
        text = read_powerpoint(file_path)
    elif file_extension == '.ipynb':
        text = read_notebook(file_path)
    elif file_extension == '.py':
        text = read_py(file_path)
    return text

def split_text(text):
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=100,
        length_function=len,
    )
    docs = text_splitter.create_documents([text])
    return docs

def custom_prompt_forall():
    return """Use the following information to answer the question in a simple, clear way.
    Rules:
    If the question is in Arabic, answer in Arabic.
    If the question is in English, answer in English.
    Prioritize the provided {context}. If the information is not in the context, you may use your own knowledge.
    Keep the answer short, clear, and easy to understand.
    Do not include reasoning steps or extra sections—only the final answer.
    {context}
    question: {question}
    answer:
    """

def initialize_chatbot():
    """Initialize the chatbot with data from the data folder"""
    try:
        data_folder = "./data"
        all_text = ""
        
        # Read all files in the data folder
        if os.path.exists(data_folder):
            for filename in os.listdir(data_folder):
                file_path = os.path.join(data_folder, filename)
                if os.path.isfile(file_path):
                    try:
                        file_content = return_content(file_path)
                        all_text += file_content + "\n\n"
                    except Exception as e:
                        print(f"Could not read {filename}: {e}")
        
        if not all_text.strip():
            # Fallback content if no files are found
            all_text = """
            Cat Care Information:
            - Cats need regular veterinary checkups
            - Provide fresh water and balanced nutrition
            - Keep litter boxes clean
            - Regular grooming is important
            - Provide scratching posts and toys
            - Spay/neuter your cats
            - Watch for signs of illness
            """
        
        # Split text and create embeddings
        docs = split_text(all_text)
        embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
        db = FAISS.from_documents(docs, embeddings)
        
        return db
        
    except Exception as e:
        return None

def get_chatbot_response(user_query, db):
    """Get response from chatbot"""
    if not db:
        return "I'm sorry, the chatbot is not available at the moment. Please try again later."
    
    try:
        # Use OpenAI model
        llm = ChatOpenAI(model_name="gpt-4o-mini", temperature=0)
        custom_prompt_template = custom_prompt_forall()
        custom_prompt = PromptTemplate(template=custom_prompt_template, input_variables=["context", "question"])
        qa_chain = RetrievalQA.from_chain_type(
            llm=llm,
            chain_type="stuff",
            retriever=db.as_retriever(),
            chain_type_kwargs={"prompt": custom_prompt}
        )
        
        response = qa_chain.invoke({"query": user_query})
        return response['result']
            
    except Exception as e:
        return f"I can help with cat care questions! Based on your query '{user_query}', I recommend consulting veterinary resources for detailed information."

chatbot_db = initialize_chatbot()